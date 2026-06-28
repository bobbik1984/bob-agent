#!/usr/bin/env python3
"""
apply_translations.py — Write LLM translations back into a .pptx file.

Applies translations at PARAGRAPH level: sets translated text on the first run
and clears remaining runs to preserve XML structure.

Usage:
    python3 apply_translations.py input.pptx translations.json output.pptx
    python3 apply_translations.py input.pptx translations.json output.pptx --backup

The translations.json should be an array of objects:
    [{"id": "s0_sp0_p0", "original": "...", "translation": "..."}, ...]

Or it can be the full extraction JSON with a "translations" key added.
"""

import argparse
import json
import re
import shutil
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip3 install python-pptx --break-system-packages", file=sys.stderr)
    sys.exit(1)


# Pattern to detect old run-level IDs (s0_sp0_p0_r0) vs new paragraph-level (s0_sp0_p0)
_RUN_ID_RE = re.compile(r"^s\d+_sp\d+_p\d+_r\d+$")
_PARA_ID_RE = re.compile(r"^s\d+_sp\d+_p\d+$")


def load_translation_map(translations_path: str) -> dict[str, str]:
    """Load translations JSON and return a dict of {element_id: translation_text}."""
    data = json.loads(Path(translations_path).read_text(encoding="utf-8"))

    # Support multiple formats:
    # 1. Direct array: [{"id": "...", "translation": "..."}]
    # 2. Extraction JSON with "translations" key added
    # 3. Extraction JSON where each element has "translation" added inline
    if isinstance(data, list):
        items = data
    elif isinstance(data, dict) and "translations" in data:
        items = data["translations"]
    elif isinstance(data, dict) and "slides" in data:
        items = []
        for slide in data["slides"]:
            for el in slide.get("elements", []):
                if "translation" in el and not el.get("skip", False):
                    items.append({
                        "id": el["element_id"],
                        "translation": el["translation"],
                        "original": el.get("text", ""),
                    })
    else:
        raise ValueError(f"Unrecognized translations format in {translations_path}")

    translation_map = {}
    for item in items:
        eid = item.get("id") or item.get("element_id")
        text = item.get("translation") or item.get("translated_text")
        if eid and text:
            translation_map[eid] = text

    return translation_map


def _detect_id_format(translation_map: dict[str, str]) -> str:
    """Detect whether translation IDs are paragraph-level or run-level."""
    for eid in translation_map:
        if _RUN_ID_RE.match(eid):
            return "run"
        if _PARA_ID_RE.match(eid):
            return "paragraph"
    return "paragraph"  # default


def iter_shapes_with_ids(shapes, slide_idx: int, prefix: str = ""):
    """Recursively yield (shape, id_prefix) pairs, including group children."""
    for shape_idx, shape in enumerate(shapes):
        shape_prefix = f"s{slide_idx}_sp{shape_idx}" if not prefix else f"{prefix}_g{shape_idx}"

        # Recurse into group shapes
        if hasattr(shape, 'shapes'):
            yield from iter_shapes_with_ids(shape.shapes, slide_idx, shape_prefix)
        elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            yield shape, shape_prefix


def apply_translations(
    input_path: str,
    translations_path: str,
    output_path: str,
    backup: bool = False,
) -> dict:
    """Apply translations and return a stats dict."""
    translation_map = load_translation_map(translations_path)
    prs = Presentation(input_path)
    id_format = _detect_id_format(translation_map)

    applied = 0
    skipped = 0
    missing = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape, shape_prefix in iter_shapes_with_ids(slide.shapes, slide_idx):
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                runs = para.runs
                if not runs:
                    continue

                if id_format == "paragraph":
                    # ── Paragraph-level application ──
                    para_id = f"{shape_prefix}_p{para_idx}"

                    # Check if all runs are empty/whitespace
                    para_text = "".join(run.text for run in runs)
                    if not para_text or not para_text.strip():
                        continue

                    if para_id in translation_map:
                        new_text = translation_map[para_id]
                        # Put full translation in the first run, clear the rest
                        runs[0].text = new_text
                        for run in runs[1:]:
                            run.text = ""
                        applied += 1
                    else:
                        skipped += 1
                        missing.append(para_id)
                else:
                    # ── Legacy run-level application (backward compat) ──
                    for run_idx, run in enumerate(runs):
                        if not run.text or not run.text.strip():
                            continue
                        element_id = f"{shape_prefix}_p{para_idx}_r{run_idx}"
                        if element_id in translation_map:
                            run.text = translation_map[element_id]
                            applied += 1
                        else:
                            skipped += 1
                            missing.append(element_id)

    # Apply notes translations if present
    notes_map = {}
    try:
        raw = json.loads(Path(translations_path).read_text(encoding="utf-8"))
        if isinstance(raw, dict) and "slides" in raw:
            for slide_data in raw["slides"]:
                idx = slide_data.get("slide_index")
                if "notes_translation" in slide_data and idx is not None:
                    notes_map[idx] = slide_data["notes_translation"]
    except Exception:
        pass

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx in notes_map:
            try:
                if slide.has_notes_slide:
                    tf = slide.notes_slide.notes_text_frame
                    if tf.paragraphs:
                        tf.paragraphs[0].runs[0].text = notes_map[slide_idx]
            except Exception:
                pass

    # Save
    if backup and Path(output_path).exists():
        shutil.copy2(output_path, output_path + ".bak")

    prs.save(output_path)

    return {
        "applied": applied,
        "skipped": skipped,
        "missing_ids": missing[:20],
        "total_missing": len(missing),
        "id_format": id_format,
    }


def main():
    parser = argparse.ArgumentParser(description="Apply translations to .pptx")
    parser.add_argument("input", help="Original .pptx file")
    parser.add_argument("translations", help="Translations JSON file")
    parser.add_argument("output", help="Output .pptx file path")
    parser.add_argument("--backup", action="store_true", help="Backup existing output file")
    args = parser.parse_args()

    for f in [args.input, args.translations]:
        if not Path(f).exists():
            print(f"ERROR: File not found: {f}", file=sys.stderr)
            sys.exit(1)

    stats = apply_translations(args.input, args.translations, args.output, backup=args.backup)

    print(f"✅ Applied {stats['applied']} translations → {args.output}", file=sys.stderr)
    print(f"   ID format: {stats['id_format']}-level", file=sys.stderr)
    if stats["total_missing"] > 0:
        print(f"⚠️  {stats['total_missing']} elements had no translation (kept original)", file=sys.stderr)
        if stats["missing_ids"]:
            print(f"   First few missing IDs: {', '.join(stats['missing_ids'][:5])}", file=sys.stderr)


if __name__ == "__main__":
    main()
