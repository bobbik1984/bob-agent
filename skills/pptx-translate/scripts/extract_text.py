#!/usr/bin/env python3
"""
extract_text.py — Extract all text from a .pptx file for translation.

Extracts at PARAGRAPH level (not run level) for coherent translation.
Each paragraph's runs are merged into one text unit.

Usage:
    python3 extract_text.py input.pptx
    python3 extract_text.py input.pptx --output extracted.json
    python3 extract_text.py input.pptx --target-lang en
"""

import argparse
import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip3 install python-pptx --break-system-packages", file=sys.stderr)
    sys.exit(1)


# ── Language detection (optional, graceful fallback) ──────────────────────────

def detect_language(texts: list[str]) -> str:
    """Detect dominant language from a sample of texts."""
    sample = " ".join(t for t in texts if len(t) > 3)[:500]
    if not sample:
        return "unknown"

    # Fast heuristic: count CJK characters
    cjk_count = sum(1 for ch in sample if "\u4e00" <= ch <= "\u9fff")
    latin_count = sum(1 for ch in sample if ch.isalpha() and ord(ch) < 128)

    if cjk_count > latin_count * 0.3:
        return "zh"

    # Try langdetect if available
    try:
        from langdetect import detect
        return detect(sample)
    except Exception:
        return "en" if latin_count > 0 else "unknown"


# ── Skip heuristics (paragraph level) ────────────────────────────────────────

_URL_RE = re.compile(r"https?://\S+|www\.\S+")
_EMAIL_RE = re.compile(r"\S+@\S+\.\S+")
_NUMBER_RE = re.compile(r"^\s*[\d\s.,/%$€¥£+\-:()]+\s*$")


def should_skip(text: str) -> tuple[bool, str]:
    """Return (should_skip, reason) for a full paragraph text."""
    stripped = text.strip()
    if not stripped:
        return True, "empty"
    if _NUMBER_RE.match(stripped):
        return True, "numbers_only"
    if _URL_RE.search(stripped):
        return True, "url"
    if _EMAIL_RE.search(stripped):
        return True, "email"
    # Only skip single-char if it's a pure symbol (not CJK, not letter)
    if len(stripped) == 1 and not stripped.isalpha() and not ("\u4e00" <= stripped <= "\u9fff"):
        return True, "single_symbol"
    return False, ""


# ── Shape type helpers ────────────────────────────────────────────────────────

def shape_type_label(shape) -> str:
    if shape.is_placeholder:
        ph = shape.placeholder_format
        ph_idx = ph.idx if ph else -1
        if ph_idx == 0:
            return "title"
        elif ph_idx == 1:
            return "subtitle"
        else:
            return "body"
    return "text_box"


# ── Main extraction ───────────────────────────────────────────────────────────

def iter_shapes_with_ids(shapes, slide_idx: int, prefix: str = ""):
    """Recursively yield (shape, id_prefix) pairs, including group children."""
    for shape_idx, shape in enumerate(shapes):
        shape_prefix = f"s{slide_idx}_sp{shape_idx}" if not prefix else f"{prefix}_g{shape_idx}"

        # Recurse into group shapes
        if hasattr(shape, 'shapes'):
            yield from iter_shapes_with_ids(shape.shapes, slide_idx, shape_prefix)
        elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            yield shape, shape_prefix


def extract(pptx_path: str, target_lang: str | None = None) -> dict:
    prs = Presentation(pptx_path)
    all_texts = []
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_elements = []

        for shape, shape_prefix in iter_shapes_with_ids(slide.shapes, slide_idx):
            shape_label = shape_type_label(shape)

            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                # Merge all runs into one paragraph text
                runs = para.runs
                if not runs:
                    continue

                para_text = "".join(run.text for run in runs)
                if not para_text or not para_text.strip():
                    continue

                element_id = f"{shape_prefix}_p{para_idx}"
                skip, skip_reason = should_skip(para_text)

                # Record run count for apply_translations to know the structure
                run_count = len(runs)

                element = {
                    "element_id": element_id,
                    "type": shape_label,
                    "text": para_text,
                    "run_count": run_count,
                    "skip": skip,
                }
                if skip:
                    element["skip_reason"] = skip_reason
                else:
                    all_texts.append(para_text)

                # Optional font metadata from first run (for context)
                try:
                    first_run = runs[0]
                    if first_run.font.size:
                        element["font_size"] = round(first_run.font.size.pt, 1)
                    if first_run.font.bold:
                        element["is_bold"] = True
                except Exception:
                    pass

                slide_elements.append(element)

        # Extract speaker notes
        notes_text = ""
        try:
            if slide.has_notes_slide:
                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    notes_text += para.text + "\n"
                notes_text = notes_text.strip()
        except Exception:
            pass

        slides_data.append({
            "slide_index": slide_idx,
            "slide_number": slide_idx + 1,
            "elements": slide_elements,
            **({  "notes": notes_text} if notes_text else {}),
        })

    detected_lang = detect_language(all_texts)

    # Determine target language
    if target_lang is None:
        if detected_lang == "zh":
            target_lang = "en"
        else:
            target_lang = "zh"

    return {
        "source_file": str(Path(pptx_path).name),
        "slide_count": len(prs.slides),
        "detected_language": detected_lang,
        "target_language": target_lang,
        "translatable_count": len(all_texts),
        "slides": slides_data,
        "all_texts": all_texts,
    }


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Extract text from .pptx for translation")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("--output", "-o", help="Output JSON file (default: stdout)")
    parser.add_argument("--target-lang", help="Target language code (e.g. en, zh, ja, fr)")
    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f"ERROR: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    result = extract(args.input, target_lang=args.target_lang)

    output_json = json.dumps(result, ensure_ascii=False, indent=2)

    if args.output:
        Path(args.output).write_text(output_json, encoding="utf-8")
        print(f"✅ Extracted {result['translatable_count']} text elements from {result['slide_count']} slides", file=sys.stderr)
        print(f"   Detected language: {result['detected_language']} → Target: {result['target_language']}", file=sys.stderr)
        print(f"   Output: {args.output}", file=sys.stderr)
    else:
        print(output_json)


if __name__ == "__main__":
    main()
