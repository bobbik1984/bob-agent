#!/usr/bin/env python3
"""
verify.py — Compare original and translated .pptx files to validate translation.

Works at PARAGRAPH level to match the extraction/application scripts.

Usage:
    python3 verify.py original.pptx translated.pptx
    python3 verify.py original.pptx translated.pptx --json
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx not installed.", file=sys.stderr)
    sys.exit(1)


def is_likely_untranslated(original: str, translated: str, source_lang: str) -> bool:
    """Heuristic: check if text appears unchanged when it should have been translated."""
    if original == translated:
        if source_lang == "zh":
            has_cjk = any("\u4e00" <= ch <= "\u9fff" for ch in original)
            return has_cjk
        elif source_lang == "en":
            latin_ratio = sum(1 for c in original if c.isalpha() and ord(c) < 128) / max(len(original), 1)
            return latin_ratio > 0.7 and len(original) > 5
    return False


def iter_shapes_with_ids(shapes, slide_idx: int, prefix: str = ""):
    """Recursively yield (shape, id_prefix) pairs, including group children."""
    for shape_idx, shape in enumerate(shapes):
        shape_prefix = f"s{slide_idx}_sp{shape_idx}" if not prefix else f"{prefix}_g{shape_idx}"
        if hasattr(shape, 'shapes'):
            yield from iter_shapes_with_ids(shape.shapes, slide_idx, shape_prefix)
        elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            yield shape, shape_prefix


def extract_all_paragraphs(pptx_path: str) -> list[dict]:
    """Extract all non-empty paragraphs with their IDs."""
    prs = Presentation(pptx_path)
    paragraphs = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape, shape_prefix in iter_shapes_with_ids(slide.shapes, slide_idx):
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                runs = para.runs
                if not runs:
                    continue
                para_text = "".join(run.text for run in runs)
                if para_text and para_text.strip():
                    paragraphs.append({
                        "id": f"{shape_prefix}_p{para_idx}",
                        "slide": slide_idx + 1,
                        "text": para_text,
                    })
    return paragraphs


def detect_lang(texts: list[str]) -> str:
    sample = " ".join(texts[:50])
    cjk = sum(1 for c in sample if "\u4e00" <= c <= "\u9fff")
    latin = sum(1 for c in sample if c.isalpha() and ord(c) < 128)
    return "zh" if cjk > latin * 0.3 else "en"


def verify(original_path: str, translated_path: str) -> dict:
    orig_paras = extract_all_paragraphs(original_path)
    trans_paras = extract_all_paragraphs(translated_path)

    issues = []
    warnings = []

    # Check slide count
    orig_prs = Presentation(original_path)
    trans_prs = Presentation(translated_path)
    orig_slide_count = len(orig_prs.slides)
    trans_slide_count = len(trans_prs.slides)

    if orig_slide_count != trans_slide_count:
        issues.append(f"Slide count mismatch: original={orig_slide_count}, translated={trans_slide_count}")

    # Check paragraph count
    if len(orig_paras) != len(trans_paras):
        warnings.append(f"Paragraph count differs: original={len(orig_paras)}, translated={len(trans_paras)}")

    # Build lookup for translated paragraphs
    trans_lookup = {p["id"]: p["text"] for p in trans_paras}
    orig_texts = [p["text"] for p in orig_paras]
    source_lang = detect_lang(orig_texts)

    untranslated = []
    for para in orig_paras:
        trans_text = trans_lookup.get(para["id"], para["text"])
        if is_likely_untranslated(para["text"], trans_text, source_lang):
            untranslated.append({
                "id": para["id"],
                "slide": para["slide"],
                "text": para["text"][:60],
            })

    if untranslated:
        warnings.append(f"{len(untranslated)} paragraphs may be untranslated (source text unchanged)")

    total_translated = len(orig_paras) - len(untranslated)
    translation_rate = round(total_translated / max(len(orig_paras), 1) * 100, 1)

    # Check for empty slides in translated version
    empty_slides = []
    for slide_idx, slide in enumerate(trans_prs.slides):
        has_text = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                has_text = True
                break
        if not has_text:
            empty_slides.append(slide_idx + 1)

    if empty_slides:
        issues.append(f"Empty slides in translated output: {empty_slides}")

    result = {
        "original": str(Path(original_path).name),
        "translated": str(Path(translated_path).name),
        "slide_count_ok": orig_slide_count == trans_slide_count,
        "original_paragraphs": len(orig_paras),
        "translated_paragraphs": len(trans_paras),
        "translation_rate_pct": translation_rate,
        "empty_slides": empty_slides,
        "issues": issues,
        "warnings": warnings,
        "untranslated_samples": untranslated[:10],
        "passed": len(issues) == 0 and translation_rate >= 80,
    }

    return result


def main():
    parser = argparse.ArgumentParser(description="Verify pptx translation completeness")
    parser.add_argument("original", help="Original .pptx")
    parser.add_argument("translated", help="Translated .pptx")
    parser.add_argument("--json", action="store_true", dest="output_json", help="Output JSON")
    args = parser.parse_args()

    for f in [args.original, args.translated]:
        if not Path(f).exists():
            print(f"ERROR: File not found: {f}", file=sys.stderr)
            sys.exit(1)

    result = verify(args.original, args.translated)

    if args.output_json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return

    # Human-readable output
    status = "✅ PASSED" if result["passed"] else "❌ ISSUES FOUND"
    print(f"\n{status} — Translation Verification")
    print(f"  Slides: {result['original_paragraphs']} paragraphs in original")
    print(f"  Translation rate: {result['translation_rate_pct']}%")

    if result["empty_slides"]:
        print(f"\n🚨 Empty slides: {result['empty_slides']}")

    if result["issues"]:
        print("\n🚨 Issues:")
        for issue in result["issues"]:
            print(f"  - {issue}")

    if result["warnings"]:
        print("\n⚠️  Warnings:")
        for w in result["warnings"]:
            print(f"  - {w}")

    if result["untranslated_samples"]:
        print(f"\n📋 Sample untranslated paragraphs (showing up to 5):")
        for item in result["untranslated_samples"][:5]:
            print(f"  Slide {item['slide']} [{item['id']}]: {item['text']!r}")

    sys.exit(0 if result["passed"] else 1)


if __name__ == "__main__":
    main()
